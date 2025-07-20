# summary_generator.py
import os
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import scipy.stats as stats
from pptx import Presentation
from pptx.util import Inches
from config import Config

class SummaryGenerator:
    def __init__(self):
        self.config = Config()
    
    def create_summary(self, folder_path):
        """総合まとめの作成"""
        datestamp = self.config.get_datestamp()
        csv_files = [f for f in os.listdir(folder_path) if f.startswith("1fData") and f.endswith(".csv")]
        
        data_list = []
        
        for file in csv_files:
            file_path_csv = os.path.join(folder_path, file)
            name_parts = file.split("_")
            if len(name_parts) >= 3:
                column_name = name_parts[1]
                df = pd.read_csv(file_path_csv)
                rent_column = [col for col in df.columns if "賃料" in col or "円" in col]
                if rent_column:
                    df_filtered = df[[rent_column[0]]]
                    df_filtered.columns = [column_name]
                    data_list.append(df_filtered)
        
        if not data_list:
            return None
        
        result_df = pd.concat(data_list, axis=1)
        stats_df = result_df.describe()
        
        # 箱ひげ図の作成
        self.create_box_plot(result_df, datestamp, folder_path)
        
        # ANOVA実施
        anova_text = self.perform_anova(result_df)
        
        # 累積比率グラフの作成
        self.create_cumulative_plot(result_df, datestamp, folder_path)
        
        # まとめのパワポ作成
        self.create_summary_powerpoint(result_df, stats_df, anova_text, datestamp, folder_path)
        
        return result_df
    
    def create_box_plot(self, result_df, datestamp, folder_path):
        """箱ひげ図の作成"""
        plt.figure(figsize=(10, 6))
        result_df.boxplot()
        plt.title(f"箱ひげ図 ({datestamp})")
        plt.ylabel("賃料（円）")
        plt.xticks(rotation=45)
        plt.grid(True)
        
        filename = f"{datestamp}_box1.png"
        image_path = os.path.join(folder_path, filename)
        plt.savefig(image_path)
        plt.close()
        
        return image_path
    
    def perform_anova(self, result_df):
        """ANOVA分析の実行"""
        result_df_clean = result_df.dropna()
        column_names = result_df_clean.columns.tolist()
        groups = [result_df_clean[col] for col in column_names]
        groups = [g for g in groups if len(g) > 0]
        
        if len(groups) > 1:
            F_value, p_value = stats.f_oneway(*groups)
            text = f"一元配置分散分析（ANOVA）の結果:\nF値: {F_value:.2f}\np値: {p_value:.3f}"
        else:
            text = "ANOVAを適用できる十分なデータがありません。"
        
        print(text)
        return text
    
    def create_cumulative_plot(self, result_df, datestamp, folder_path):
        """累積比率グラフの作成"""
        plt.figure(figsize=(10, 6))
        
        for column in result_df.columns:
            data = np.sort(result_df[column].dropna())
            cum_data = np.cumsum(data) / np.sum(data)
            plt.plot(data, cum_data, label=column)
        
        plt.xlabel("賃料（円）")
        plt.ylabel("累積賃料比率")
        plt.title(f"賃料の累積比率グラフ ({datestamp})")
        plt.legend()
        plt.grid(True)
        
        filename = f"{datestamp}_cum1.png"
        image_path = os.path.join(folder_path, filename)
        plt.savefig(image_path)
        plt.close()
        
        return image_path
    
    def create_summary_powerpoint(self, result_df, stats_df, anova_text, datestamp, folder_path):
        """まとめパワーポイントの作成"""
        ppt = Presentation()
        
        # タイトルスライド
        slide_layout = ppt.slide_layouts[0]
        slide = ppt.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = f"データサマリー ({datestamp})"
        subtitle = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
        subtitle.text = "各駅の賃料をまとめました"
        
        # 基礎統計量スライド
        slide_layout = ppt.slide_layouts[5]
        slide = ppt.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = "基礎統計量(小数桁数ご容赦)"
        
        rows, cols = stats_df.shape
        table = slide.shapes.add_table(rows+1, cols+1, Inches(1), Inches(1.5), Inches(8), Inches(4)).table
        
        table.cell(0, 0).text = "統計項目"
        for col_idx, col_name in enumerate(stats_df.columns):
            table.cell(0, col_idx+1).text = col_name
        
        for row_idx, (index, row_data) in enumerate(stats_df.iterrows()):
            table.cell(row_idx+1, 0).text = index
            for col_idx, value in enumerate(row_data):
                table.cell(row_idx+1, col_idx+1).text = f"{value:.2f}"
        
        # 箱ひげ図スライド
        slide_layout = ppt.slide_layouts[1]
        slide = ppt.slides.add_slide(slide_layout)
        slide.shapes.title.text = "賃料の箱ひげ図"
        
        box_image_path = os.path.join(folder_path, f"{datestamp}_box1.png")
        if os.path.exists(box_image_path):
            slide.shapes.add_picture(box_image_path, Inches(1), Inches(2), Inches(8), Inches(5))
        
        text_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(5), Inches(1.2))
        text_frame = text_box.text_frame
        text_frame.text = anova_text
        text_frame.word_wrap = True
        
        # 累積比率グラフスライド
        slide_layout = ppt.slide_layouts[1]
        slide = ppt.slides.add_slide(slide_layout)
        slide.shapes.title.text = "賃料の分布"
        
        cum_image_path = os.path.join(folder_path, f"{datestamp}_cum1.png")
        if os.path.exists(cum_image_path):
            slide.shapes.add_picture(cum_image_path, Inches(1), Inches(2), Inches(8), Inches(5))
        
        # PowerPointファイルの保存
        ppt_filename = f"1c_{datestamp}_sum.pptx"
        ppt_path = os.path.join(folder_path, ppt_filename)
        ppt.save(ppt_path)
        print(f"総合まとめPowerPointを保存しました: {ppt_path}")
        
        return ppt_path