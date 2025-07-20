# powerpoint_generator.py
import os
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Cm, Pt
from pptx.dml.color import RGBColor
from config import Config

class PowerPointGenerator:
    def __init__(self):
        self.config = Config()
    
    def create_presentation(self, station, n, timestamp, folder_path, df_base1, cat1, 
                          df_stats11, df_stats12, df_mrl1, df_mrl2, df_vif1, df_comp1):
        """プレゼンテーション作成"""
        ppt = Presentation()
        
        self.add_title_slide(ppt, station, timestamp, n)
        self.add_basic_info_slide(ppt, df_base1)
        self.add_category_slide(ppt, cat1)
        self.add_stats_slides(ppt, df_stats11, df_stats12)
        self.add_graph_slides(ppt, station, folder_path)
        self.add_regression_slides(ppt, df_mrl1, df_mrl2, df_vif1, df_comp1, station, folder_path)
        self.add_footers(ppt, station, n, timestamp)
        
        return ppt
    
    def add_title_slide(self, ppt, station, timestamp, n):
        """タイトルスライドの追加"""
        slide_layout = ppt.slide_layouts[0]
        slide = ppt.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = f"{station}駅\n徒歩圏内の賃貸物件の\n調査結果"
        subtitle.text = f"調査時刻: {timestamp}\nデータ件数は{n}です\n ご注意:重複はなるべく排除していますが排除され切れていません"
    
    def add_basic_info_slide(self, ppt, df_base1):
        """基本情報スライドの追加"""
        slide_layout = ppt.slide_layouts[6]
        slide = ppt.slides.add_slide(slide_layout)
        
        text_box = slide.shapes.add_textbox(Cm(0.4), Cm(0.4), Cm(5), Cm(1))
        text_frame = text_box.text_frame
        p = text_frame.add_paragraph()
        p.text = "基本情報"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 0, 0)
        
        self.add_dataframe_table(slide, df_base1, Cm(1.5), Cm(2), Cm(22), Cm(4))
    
    def add_category_slide(self, ppt, cat1):
        """カテゴリー情報スライドの追加"""
        slide_layout = ppt.slide_layouts[6]
        slide = ppt.slides.add_slide(slide_layout)
        
        text_box = slide.shapes.add_textbox(Cm(0.4), Cm(0.4), Cm(5), Cm(1))
        text_frame = text_box.text_frame
        p = text_frame.add_paragraph()
        p.text = "カテゴリー情報"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 0, 0)
        
        self.add_dataframe_table(slide, cat1, Cm(1.5), Cm(2), Cm(22), Cm(15))
    
    def add_stats_slides(self, ppt, df_stats11, df_stats12):
        """統計スライドの追加"""
        # 統計A
        slide_layout = ppt.slide_layouts[6]
        slide = ppt.slides.add_slide(slide_layout)
        
        text_box = slide.shapes.add_textbox(Cm(0.4), Cm(0.4), Cm(5), Cm(1))
        text_frame = text_box.text_frame
        p = text_frame.add_paragraph()
        p.text = "基礎統計量情報A"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 0, 0)
        
        self.add_dataframe_table(slide, df_stats11, Cm(1.5), Cm(2), Cm(22), Cm(15))
        
        # 統計B
        slide_layout = ppt.slide_layouts[6]
        slide = ppt.slides.add_slide(slide_layout)
        
        text_box = slide.shapes.add_textbox(Cm(0.4), Cm(0.4), Cm(5), Cm(1))
        text_frame = text_box.text_frame
        p = text_frame.add_paragraph()
        p.text = "基礎統計量情報B"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 0, 0)
        
        self.add_dataframe_table(slide, df_stats12, Cm(1.5), Cm(2), Cm(22), Cm(15))
    
    def add_graph_slides(self, ppt, station, folder_path):
        """グラフスライドの追加"""
        datestamp = self.config.get_datestamp()
        
        # 分布と一次回帰のグラフスライド
        slide_layout = ppt.slide_layouts[3]
        slide = ppt.slides.add_slide(slide_layout)
        
        text_boxes = [shape for shape in slide.shapes if shape.has_text_frame]
        
        if len(text_boxes) >= 2:
            text_boxes[0].text = "全体の分布と一次回帰のグラフ"
            text_boxes[1].text = "分布"
            text_boxes[2].text = "一次回帰"
        
        tg1 = os.path.normpath(os.path.join(folder_path, f"{station}_{datestamp}_tg1.png"))
        tg2 = os.path.normpath(os.path.join(folder_path, f"{station}_{datestamp}_tg2.png"))
        
        if os.path.exists(tg1) and os.path.exists(tg2) and len(text_boxes) >= 3:
            left_placeholder = text_boxes[1]
            right_placeholder = text_boxes[2]
            
            slide.shapes.add_picture(tg1, left_placeholder.left, left_placeholder.top, 
                                   left_placeholder.width, left_placeholder.height)
            slide.shapes.add_picture(tg2, right_placeholder.left, right_placeholder.top, 
                                   right_placeholder.width, right_placeholder.height)
        
        # 個別グラフスライド
        graph_configs = [
            ("賃料分布グラフ", f"{station}_{datestamp}_gr1.png"),
            ("徒歩時間グラフ", f"{station}_{datestamp}_gw1.png"),
            ("専有面積グラフ", f"{station}_{datestamp}_gs1.png"),
            ("築年数グラフ", f"{station}_{datestamp}_ga1.png")
        ]
        
        for title, filename in graph_configs:
            self.add_image_slide(ppt, title, os.path.join(folder_path, filename))
        
        # 散布図スライド
        scatter_configs = [
            ("賃料と徒歩時間の散布図", f"{station}_{datestamp}_tgscat1.png"),
            ("賃料と専有面積の散布図", f"{station}_{datestamp}_tgscat2.png"),
            ("賃料と築年数の散布図", f"{station}_{datestamp}_tgscat3.png")
        ]
        
        for title, filename in scatter_configs:
            self.add_image_slide(ppt, title, os.path.join(folder_path, filename))
    
    def add_regression_slides(self, ppt, df_mrl1, df_mrl2, df_vif1, df_comp1, station, folder_path):
        """回帰分析スライドの追加"""
        datestamp = self.config.get_datestamp()
        
        # 重回帰分析結果スライド
        slide_layout = ppt.slide_layouts[5]
        slide = ppt.slides.add_slide(slide_layout)
        
        if slide.shapes.title:
            slide.shapes.title.text = "重回帰分析結果"
        
        text_box = slide.shapes.add_textbox(Cm(0.4), Cm(2), Cm(5), Cm(1))
        text_frame = text_box.text_frame
        p = text_frame.add_paragraph()
        p.text = "重回帰基礎結果とcoefficients"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 0, 0)
        
        # テーブル追加
        self.add_dataframe_table_custom(slide, df_mrl1, Cm(1.5), Cm(4.0), Cm(22), Cm(4))
        self.add_dataframe_table_custom(slide, df_mrl2, Cm(1.5), Cm(10.0), Cm(22), Cm(4))
        
        # VIFスライド
        slide_layout = ppt.slide_layouts[6]
        slide = ppt.slides.add_slide(slide_layout)
        
        text_box = slide.shapes.add_textbox(Cm(0.4), Cm(0.5), Cm(5), Cm(1))
        text_frame = text_box.text_frame
        p = text_frame.add_paragraph()
        p.text = "重回帰の多重共線性（VIF)"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 0, 0)
        
        self.add_dataframe_table(slide, df_vif1, Cm(1.5), Cm(2.5), Cm(22), Cm(4))
        
        # 予測家賃と実家賃の関係スライド
        image_path = os.path.join(folder_path, f"{station}_{datestamp}_mlrap1.png")
        self.add_image_slide(ppt, "予測家賃と実家賃の関係", image_path)
        
        # 面積毎の家賃予測スライド
        slide_layout = ppt.slide_layouts[6]
        slide = ppt.slides.add_slide(slide_layout)
        
        text_box = slide.shapes.add_textbox(Cm(0.4), Cm(0.5), Cm(5), Cm(1))
        text_frame = text_box.text_frame
        p = text_frame.add_paragraph()
        p.text = "面積毎の家賃予測"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 0, 0)
        
        self.add_dataframe_table(slide, df_comp1, Cm(1.5), Cm(2), Cm(22), Cm(4))
    
    def add_image_slide(self, ppt, title, image_path):
        """画像スライドの追加"""
        slide_layout = ppt.slide_layouts[5]
        slide = ppt.slides.add_slide(slide_layout)
        
        if slide.shapes.title:
            slide.shapes.title.text = title
        
        if os.path.exists(image_path):
            left = Inches(0.3)
            top = Inches(1.5)
            width = Inches(9.5)
            height = Inches(4.5)
            slide.shapes.add_picture(image_path, left, top, width, height)
    
    def add_dataframe_table(self, slide, df, left, top, width, height):
        """データフレームをテーブルとして追加"""
        rows, cols = df.shape[0] + 1, df.shape[1]
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        
        # ヘッダー行の設定
        for col_idx, col_name in enumerate(df.columns):
            cell = table.cell(0, col_idx)
            cell.text = str(col_name)
        
        # データ行の設定
        for row_idx, row in enumerate(df.itertuples(), start=1):
            for col_idx, value in enumerate(row[1:]):
                cell = table.cell(row_idx, col_idx)
                cell.text = str(value)
    
    def add_dataframe_table_custom(self, slide, df, left, top, width, height):
        """カスタムデータフレームテーブル追加"""
        table = slide.shapes.add_table(df.shape[0], df.shape[1], left, top, width, height).table
        
        for row_idx, (index, row) in enumerate(df.iterrows()):
            for col_idx, value in enumerate(row):
                cell = table.cell(row_idx, col_idx)
                cell.text = str(value)
    
    def add_footers(self, ppt, station, n, timestamp):
        """フッターの追加"""
        for index, slide in enumerate(ppt.slides):
            current_page = index + 1
            total_pages = len(ppt.slides)
            
            left_text = f"{station}, n={n}"
            center_text = f"{current_page}/{total_pages}"
            right_text = f"{timestamp}"
            
            # テキストボックスの追加
            left_box = slide.shapes.add_textbox(Inches(0.4), Inches(7.15), Inches(2), Inches(0.3))
            left_box.text_frame.text = left_text
            
            center_box = slide.shapes.add_textbox(Inches(4.2), Inches(7.15), Inches(2), Inches(0.3))
            center_box.text_frame.text = center_text
            
            right_box = slide.shapes.add_textbox(Inches(8.0), Inches(7.15), Inches(2), Inches(0.3))
            right_box.text_frame.text = right_text